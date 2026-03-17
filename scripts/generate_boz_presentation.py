"""
Generate a comprehensive Bank of Zambia presentation for Monde
Payment System Business Designation Application

Usage: python scripts/generate_boz_presentation.py
Output: Monde_BoZ_Presentation.pptx in project root
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Brand colors ──────────────────────────────────────────────
MONDE_BLUE    = RGBColor(0x1A, 0x56, 0xDB)
MONDE_DARK    = RGBColor(0x0F, 0x17, 0x2A)
MONDE_GREEN   = RGBColor(0x10, 0xB9, 0x81)
MONDE_GOLD    = RGBColor(0xF5, 0xA6, 0x23)
MONDE_RED     = RGBColor(0xEF, 0x44, 0x44)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF1, 0xF5, 0xF9)
MEDIUM_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT      = RGBColor(0x1E, 0x29, 0x3B)
SUBTITLE_GRAY = RGBColor(0x64, 0x74, 0x8B)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    # line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=16,
                          default_color=DARK_TEXT, font_name='Calibri', line_spacing=1.15):
    """lines = list of (text, font_size, color, bold, alignment) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
        # line spacing
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return txBox

def add_bullet_textbox(slide, left, top, width, height, title, bullets,
                       title_size=20, bullet_size=15, title_color=DARK_TEXT,
                       bullet_color=SUBTITLE_GRAY, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = font_name
    p.space_after = Pt(8)
    # Bullets
    for bullet in bullets:
        bp = tf.add_paragraph()
        bp.text = bullet
        bp.font.size = Pt(bullet_size)
        bp.font.color.rgb = bullet_color
        bp.font.name = font_name
        bp.space_before = Pt(4)
        bp.space_after = Pt(2)
        bp.level = 0
        # bullet character
        pPr = bp._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        pPr.append(buChar)
        buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})
        pPr.append(buSzPct)
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': 'F5A623'})
        buClr.append(srgb)
        pPr.append(buClr)
        # indent
        pPr.set('marL', str(Pt(20)))
        pPr.set('indent', str(Pt(-14)))
    return txBox

def add_icon_card(slide, left, top, width, height, icon_text, title, subtitle,
                  bg_color=WHITE, accent_color=MONDE_BLUE, border=True):
    """Add a card with an icon circle, title and subtitle"""
    card = add_shape(slide, left, top, width, height, bg_color, corner_radius=0.05)
    if border:
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1)
    # Icon circle
    circle_size = Inches(0.55)
    cx = left + Inches(0.25)
    cy = top + Inches(0.25)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    # Icon text
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.95), width - Inches(0.5), Inches(0.35),
                title, font_size=14, color=DARK_TEXT, bold=True)
    # Subtitle
    add_textbox(slide, left + Inches(0.25), top + Inches(1.25), width - Inches(0.5), Inches(0.8),
                subtitle, font_size=11, color=SUBTITLE_GRAY)
    return card

def add_status_row(slide, left, top, width, label, status, detail=""):
    """Add a compliance status row with colored indicator"""
    row_height = Inches(0.45)
    # Background
    bg = add_shape(slide, left, top, width, row_height, LIGHT_GRAY, corner_radius=0.02)
    # Status dot
    dot_size = Inches(0.15)
    if status == "compliant":
        dot_color = MONDE_GREEN
    elif status == "partial":
        dot_color = MONDE_GOLD
    else:
        dot_color = MONDE_RED
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15),
                                  top + Inches(0.15), dot_size, dot_size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    # Label
    add_textbox(slide, left + Inches(0.4), top + Inches(0.07), width * 0.5, row_height,
                label, font_size=12, color=DARK_TEXT, bold=True)
    # Detail
    if detail:
        add_textbox(slide, left + width * 0.55, top + Inches(0.07), width * 0.42, row_height,
                    detail, font_size=11, color=SUBTITLE_GRAY)


def add_slide_number(slide, num, total):
    add_textbox(slide, SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.5),
                Inches(1.3), Inches(0.3), f"{num} / {total}",
                font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header_accent(slide):
    """Add the gold accent bar at top-left"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MONDE_GOLD
    bar.line.fill.background()
    return bar

# ══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

TOTAL_SLIDES = 18

def slide_01_title(prs):
    """Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, MONDE_DARK)
    # Gold accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.0), Inches(2.0), Inches(0.08), Inches(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = MONDE_GOLD
    line.line.fill.background()
    # Company name
    add_textbox(slide, Inches(1.4), Inches(1.8), Inches(8), Inches(0.8),
                "MONDE", font_size=52, color=WHITE, bold=True)
    # Tagline
    add_textbox(slide, Inches(1.4), Inches(2.5), Inches(8), Inches(0.5),
                "Tap. Pay. Done.", font_size=24, color=MONDE_GOLD, bold=False)
    # Main title
    add_textbox(slide, Inches(1.4), Inches(3.4), Inches(10), Inches(1.0),
                "Payment System Business\nDesignation Application",
                font_size=36, color=WHITE, bold=True, line_spacing=1.2)
    # Subtitle
    add_textbox(slide, Inches(1.4), Inches(4.8), Inches(10), Inches(0.5),
                "Presentation to the Bank of Zambia", font_size=18, color=MEDIUM_GRAY)
    # Date
    add_textbox(slide, Inches(1.4), Inches(5.4), Inches(10), Inches(0.4),
                "National Payment Systems Act No. 1 of 2007  |  Section 12", font_size=14, color=SUBTITLE_GRAY)
    # Bottom bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MONDE_BLUE
    bar.line.fill.background()
    add_slide_number(slide, 1, TOTAL_SLIDES)

def slide_02_agenda(prs):
    """Agenda / Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.6),
                "AGENDA", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Presentation Overview", font_size=32, color=DARK_TEXT, bold=True)

    items = [
        ("01", "Executive Summary", "Who we are and what we seek"),
        ("02", "The Problem", "Financial inclusion gap in Zambia"),
        ("03", "Our Solution — Monde", "Platform overview and key features"),
        ("04", "Technology & Security", "Architecture, data protection, fraud prevention"),
        ("05", "Agent Network Model", "Cash-in, cash-out, and agent operations"),
        ("06", "Business Model & Fees", "Revenue structure and fee transparency"),
        ("07", "Regulatory Compliance", "NPSA 2007 mapping and compliance assessment"),
        ("08", "Compliance Roadmap", "Steps toward full designation readiness"),
        ("09", "Risk Management & AML/CFT", "Controls, limits, and monitoring"),
        ("10", "Market Opportunity", "Zambia's digital payments landscape"),
        ("11", "Implementation Timeline", "Phased rollout plan"),
        ("12", "Next Steps & Contact", "Path forward with the Bank of Zambia"),
    ]
    col1_x = Inches(1.2)
    col2_x = Inches(7.0)
    y_start = Inches(2.0)
    row_h = Inches(0.42)
    for i, (num, title, desc) in enumerate(items):
        col = col1_x if i < 6 else col2_x
        y = y_start + (i % 6) * row_h
        add_textbox(slide, col, y, Inches(0.5), row_h,
                    num, font_size=14, color=MONDE_BLUE, bold=True)
        add_textbox(slide, col + Inches(0.5), y, Inches(2.5), row_h,
                    title, font_size=13, color=DARK_TEXT, bold=True)
        add_textbox(slide, col + Inches(3.1), y, Inches(2.5), row_h,
                    desc, font_size=11, color=SUBTITLE_GRAY)
    add_slide_number(slide, 2, TOTAL_SLIDES)

def slide_03_exec_summary(prs):
    """Executive Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "01  |  EXECUTIVE SUMMARY", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Monde at a Glance", font_size=32, color=DARK_TEXT, bold=True)

    # Left column — description
    desc = (
        "Monde is a mobile payment application purpose-built for Zambia's "
        "financial ecosystem. Our platform enables instant, secure, and affordable "
        "person-to-person (P2P) payments, QR code transactions, NFC tap-to-pay, "
        "and agent-assisted cash-in/cash-out services.\n\n"
        "We are seeking designation as a Payment System Business under Section 12 "
        "of the National Payment Systems Act No. 1 of 2007, to legally operate "
        "money transfer and transmission services within the Republic of Zambia.\n\n"
        "Monde is designed to be fully interoperable across all major Zambian "
        "payment providers — Airtel Money, MTN MoMo, Zamtel Kwacha, FNB, Zanaco, "
        "and Absa — bridging the gap between mobile money ecosystems and traditional banking."
    )
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(5.5), Inches(4.5),
                desc, font_size=14, color=SUBTITLE_GRAY, line_spacing=1.4)

    # Right column — key stats cards
    cards = [
        ("6", "Providers\nSupported", MONDE_BLUE),
        ("3", "Max Taps Per\nTransaction", MONDE_GREEN),
        ("ZMW", "Zambian\nKwacha", MONDE_GOLD),
        ("24/7", "Real-Time\nSettlement", RGBColor(0x8B, 0x5C, 0xF6)),
    ]
    card_w = Inches(2.5)
    card_h = Inches(1.2)
    for i, (num, label, color) in enumerate(cards):
        cx = Inches(7.5) + (i % 2) * (card_w + Inches(0.2))
        cy = Inches(1.8) + (i // 2) * (card_h + Inches(0.2))
        card = add_shape(slide, cx, cy, card_w, card_h, LIGHT_GRAY, corner_radius=0.04)
        add_textbox(slide, cx + Inches(0.25), cy + Inches(0.15), card_w - Inches(0.3), Inches(0.5),
                    num, font_size=28, color=color, bold=True)
        add_textbox(slide, cx + Inches(0.25), cy + Inches(0.7), card_w - Inches(0.3), Inches(0.5),
                    label, font_size=11, color=SUBTITLE_GRAY)
    add_slide_number(slide, 3, TOTAL_SLIDES)

def slide_04_problem(prs):
    """The Problem — Financial Inclusion Gap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MONDE_DARK)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "02  |  THE PROBLEM", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(11), Inches(0.8),
                "Zambia's Payment Ecosystem Challenges",
                font_size=32, color=WHITE, bold=True)

    challenges = [
        ("Fragmented Ecosystem",
         "Multiple mobile money providers (Airtel, MTN, Zamtel) and banks (FNB, Zanaco, Absa) "
         "operate in silos. Sending money across providers is cumbersome, expensive, or impossible.",
         MONDE_BLUE),
        ("Low Interoperability",
         "Despite BoZ's National Financial Switch (NFS) initiative, cross-platform P2P payments "
         "remain limited. Users are forced to maintain multiple accounts.",
         MONDE_GOLD),
        ("High Transaction Costs",
         "Cross-provider transfers incur high fees that disproportionately burden low-income users. "
         "Small transactions become uneconomical.",
         MONDE_RED),
        ("Cash Dependency",
         "Zambia remains heavily cash-dependent. The unbanked and underbanked need accessible "
         "digital onramps through agent networks and simple mobile interfaces.",
         MONDE_GREEN),
    ]
    card_w = Inches(5.3)
    card_h = Inches(1.5)
    for i, (title, desc, color) in enumerate(challenges):
        cx = Inches(1.2) + (i % 2) * (card_w + Inches(0.3))
        cy = Inches(2.2) + (i // 2) * (card_h + Inches(0.3))
        card = add_shape(slide, cx, cy, card_w, card_h, RGBColor(0x1A, 0x20, 0x35), corner_radius=0.04)
        card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
        card.line.width = Pt(1)
        # colored bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.06), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, cx + Inches(0.3), cy + Inches(0.15), card_w - Inches(0.5), Inches(0.35),
                    title, font_size=16, color=WHITE, bold=True)
        add_textbox(slide, cx + Inches(0.3), cy + Inches(0.55), card_w - Inches(0.5), Inches(0.9),
                    desc, font_size=11, color=MEDIUM_GRAY, line_spacing=1.3)
    add_slide_number(slide, 4, TOTAL_SLIDES)

def slide_05_solution(prs):
    """Our Solution — Monde Platform"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "03  |  OUR SOLUTION", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Monde — Universal Payment Platform", font_size=32, color=DARK_TEXT, bold=True)

    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(0.5),
                "Every transaction completes in 3 taps or fewer. Interoperable across all Zambian providers.",
                font_size=14, color=SUBTITLE_GRAY)

    features = [
        ("QR", "QR Code Payments",
         "Scan and pay instantly at merchants or between individuals. Dynamic and static QR codes supported.",
         MONDE_BLUE),
        ("NFC", "Tap to Pay",
         "Hold phones together to transfer money. Near-field communication for instant proximity payments.",
         MONDE_GREEN),
        ("P2P", "Send & Receive",
         "Send money to any Zambian phone number across Airtel, MTN, Zamtel, FNB, Zanaco, Absa.",
         MONDE_GOLD),
        ("W", "Digital Wallet",
         "Top up from mobile money, manage balance, withdraw to any linked account. Real-time updates.",
         RGBColor(0x8B, 0x5C, 0xF6)),
        ("AG", "Agent Network",
         "Cash-in/cash-out through designated agents. Commission-based model with anti-fraud controls.",
         RGBColor(0xEC, 0x48, 0x99)),
        ("TX", "Transaction History",
         "Full audit trail with date grouping, filters, receipt generation, and share functionality.",
         RGBColor(0x06, 0xB6, 0xD4)),
    ]
    card_w = Inches(3.5)
    card_h = Inches(2.0)
    for i, (icon, title, desc, color) in enumerate(features):
        cx = Inches(1.0) + (i % 3) * (card_w + Inches(0.25))
        cy = Inches(2.6) + (i // 3) * (card_h + Inches(0.25))
        add_icon_card(slide, cx, cy, card_w, card_h, icon, title, desc,
                     accent_color=color)
    add_slide_number(slide, 5, TOTAL_SLIDES)

def slide_06_tech_security(prs):
    """Technology & Security Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "04  |  TECHNOLOGY & SECURITY", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Enterprise-Grade Security Architecture", font_size=32, color=DARK_TEXT, bold=True)

    # Left: Tech stack
    add_bullet_textbox(slide, Inches(1.0), Inches(1.9), Inches(5.5), Inches(5.0),
                       "Technology Stack", [
                           "React Native + Expo SDK 52 (cross-platform mobile)",
                           "TypeScript for type-safe, maintainable codebase",
                           "Supabase (PostgreSQL) — auth, database, real-time subscriptions",
                           "Lipila API — payment gateway for MoMo collections/disbursements",
                           "Row-Level Security (RLS) on all database tables",
                           "Edge Functions for server-side payment processing",
                           "Real-time balance and transaction updates via WebSockets",
                       ])

    # Right: Security
    sec_items = [
        ("PIN Auth", "4-digit PIN with secure hashing.\nNo plain-text storage.", MONDE_BLUE),
        ("Auto-Lock", "App locks after 2 minutes of\ninactivity. Full PIN re-entry.", MONDE_GREEN),
        ("Rate Limiting", "5 login attempts max.\n30-second lockout period.", MONDE_GOLD),
        ("Input Sanitization", "All inputs sanitized against\ninjection attacks.", RGBColor(0x8B, 0x5C, 0xF6)),
        ("Amount Controls", "Min K1, max K50,000.\nBalance validation pre-transaction.", MONDE_RED),
        ("RLS Policies", "Database row-level security.\nUsers only access own data.", RGBColor(0x06, 0xB6, 0xD4)),
    ]
    card_w = Inches(2.35)
    card_h = Inches(1.15)
    for i, (title, desc, color) in enumerate(sec_items):
        cx = Inches(7.0) + (i % 2) * (card_w + Inches(0.15))
        cy = Inches(1.9) + (i // 2) * (card_h + Inches(0.15))
        card = add_shape(slide, cx, cy, card_w, card_h, LIGHT_GRAY, corner_radius=0.03)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.05), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.1), card_w - Inches(0.3), Inches(0.25),
                    title, font_size=12, color=DARK_TEXT, bold=True)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.4), card_w - Inches(0.3), Inches(0.7),
                    desc, font_size=10, color=SUBTITLE_GRAY, line_spacing=1.3)
    add_slide_number(slide, 6, TOTAL_SLIDES)

def slide_07_data_protection(prs):
    """Data Protection & Privacy"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MONDE_DARK)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "04B  |  DATA PROTECTION", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Data Protection Act No. 3 of 2021 Compliance",
                font_size=32, color=WHITE, bold=True)

    left_items = [
        ("Data Minimization",
         "Only essential user data collected: phone number, name, transaction records. "
         "No unnecessary personal data harvested."),
        ("Secure Storage",
         "All data stored in Supabase PostgreSQL with encryption at rest. "
         "Row-Level Security ensures users can only access their own records."),
        ("Consent-Based Processing",
         "User data processed only with explicit consent at registration. "
         "Clear terms of service presented before account creation."),
        ("Right to Access & Deletion",
         "Users can view all their data through the app. Account deletion "
         "procedures will be implemented as part of compliance readiness."),
    ]
    right_items = [
        ("Confidential Transmission",
         "All API calls over HTTPS/TLS. Supabase enforces encrypted "
         "connections. No data transmitted in plain text."),
        ("Access Controls",
         "JWT-based authentication with session refresh. Admin access "
         "restricted to designated admin accounts only."),
        ("Audit Trail",
         "Complete transaction history maintained. All payment operations "
         "logged with timestamps, references, and status tracking."),
        ("Data Retention",
         "Records maintained for minimum 6 years per NPSA Section 28. "
         "Electronic storage with integrity verification."),
    ]
    card_w = Inches(5.3)
    card_h = Inches(1.0)
    for i, (title, desc) in enumerate(left_items):
        cy = Inches(2.0) + i * (card_h + Inches(0.15))
        card = add_shape(slide, Inches(1.0), cy, card_w, card_h,
                        RGBColor(0x1A, 0x20, 0x35), corner_radius=0.03)
        card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
        card.line.width = Pt(1)
        add_textbox(slide, Inches(1.25), cy + Inches(0.08), card_w - Inches(0.4), Inches(0.25),
                    title, font_size=13, color=MONDE_GREEN, bold=True)
        add_textbox(slide, Inches(1.25), cy + Inches(0.38), card_w - Inches(0.4), Inches(0.55),
                    desc, font_size=10, color=MEDIUM_GRAY, line_spacing=1.3)
    for i, (title, desc) in enumerate(right_items):
        cy = Inches(2.0) + i * (card_h + Inches(0.15))
        card = add_shape(slide, Inches(6.8), cy, card_w, card_h,
                        RGBColor(0x1A, 0x20, 0x35), corner_radius=0.03)
        card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
        card.line.width = Pt(1)
        add_textbox(slide, Inches(7.05), cy + Inches(0.08), card_w - Inches(0.4), Inches(0.25),
                    title, font_size=13, color=MONDE_BLUE, bold=True)
        add_textbox(slide, Inches(7.05), cy + Inches(0.38), card_w - Inches(0.4), Inches(0.55),
                    desc, font_size=10, color=MEDIUM_GRAY, line_spacing=1.3)
    add_slide_number(slide, 7, TOTAL_SLIDES)

def slide_08_agent_network(prs):
    """Agent Network Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "05  |  AGENT NETWORK", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Agent Cash-In / Cash-Out Model", font_size=32, color=DARK_TEXT, bold=True)

    # Agent operations table
    operations = [
        ("Cash-In (Deposit)", "Agent accepts cash from customer, credits their Monde wallet digitally",
         "0.5% commission to agent (Monde pays)", MONDE_GREEN),
        ("Cash-Out", "Customer generates 6-digit code, agent verifies and gives cash",
         "Tiered fees: K2.50–K50 (70/30 agent/Monde split)", MONDE_BLUE),
        ("Agent Transfer", "Agent-to-agent float transfer for liquidity management",
         "Free, K50,000 daily cap", MONDE_GOLD),
        ("Float Top-Up", "Agent loads wallet balance via mobile money",
         "Standard top-up fees apply", RGBColor(0x8B, 0x5C, 0xF6)),
    ]
    y = Inches(2.0)
    for title, desc, fee, color in operations:
        row = add_shape(slide, Inches(1.0), y, Inches(11.2), Inches(0.85), LIGHT_GRAY, corner_radius=0.02)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, Inches(0.06), Inches(0.85))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, Inches(1.3), y + Inches(0.08), Inches(2.5), Inches(0.3),
                    title, font_size=14, color=DARK_TEXT, bold=True)
        add_textbox(slide, Inches(1.3), y + Inches(0.42), Inches(4.5), Inches(0.35),
                    desc, font_size=11, color=SUBTITLE_GRAY)
        add_textbox(slide, Inches(7.5), y + Inches(0.2), Inches(4.5), Inches(0.4),
                    fee, font_size=12, color=color, bold=True)
        y += Inches(1.0)

    # Anti-fraud section
    add_textbox(slide, Inches(1.2), Inches(6.0), Inches(4), Inches(0.4),
                "Anti-Fraud Controls", font_size=18, color=DARK_TEXT, bold=True)
    fraud_items = [
        "Circular fraud block — agent cannot cash-out for customer they deposited to within 24h",
        "Daily deposit limit — max 3 deposits per customer per day from any agent",
        "Agent transfer daily cap — K50,000/day per agent",
        "Self-service blocks — agents cannot deposit to self or cash-out own requests",
        "Admin-only agent designation — only admin can toggle agent status",
    ]
    for i, item in enumerate(fraud_items):
        add_textbox(slide, Inches(1.4), Inches(6.45) + i * Inches(0.22), Inches(10), Inches(0.22),
                    f"\u2022  {item}", font_size=10, color=SUBTITLE_GRAY)
    add_slide_number(slide, 8, TOTAL_SLIDES)

def slide_09_business_model(prs):
    """Business Model & Fee Structure"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "06  |  BUSINESS MODEL", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Revenue Model & Fee Transparency", font_size=32, color=DARK_TEXT, bold=True)

    # Fee table
    headers = ["Service", "Fee to Customer", "Monde Revenue", "Provider Cost"]
    rows = [
        ["P2P Send", "1.5% of amount", "1.5% (less provider cost)", "Included"],
        ["Top-Up (MoMo)", "2.5% collection fee", "Margin on fee", "2.5% (Lipila)"],
        ["Withdrawal (MoMo)", "1.5% disbursement fee", "Margin on fee", "1.5% (Lipila)"],
        ["Cash-In (Agent)", "Free to customer", "Agent commission from Monde", "0.5% to agent"],
        ["Cash-Out (Agent)", "K2.50 – K50 tiered", "30% of fee", "70% to agent"],
        ["Agent Transfer", "Free", "—", "—"],
        ["QR/NFC Payment", "Same as P2P Send", "1.5%", "Included"],
    ]

    table_left = Inches(1.0)
    table_top = Inches(2.0)
    col_widths = [Inches(2.5), Inches(2.8), Inches(2.8), Inches(2.8)]
    row_height = Inches(0.45)

    # Header row
    x = table_left
    for j, header in enumerate(headers):
        cell_bg = add_shape(slide, x, table_top, col_widths[j], row_height, MONDE_BLUE, corner_radius=0.0)
        cell_bg.adjustments[0] = 0
        add_textbox(slide, x + Inches(0.15), table_top + Inches(0.07), col_widths[j] - Inches(0.2), row_height,
                    header, font_size=12, color=WHITE, bold=True)
        x += col_widths[j]

    # Data rows
    for i, row in enumerate(rows):
        y = table_top + (i + 1) * row_height
        bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
        x = table_left
        for j, cell in enumerate(row):
            cell_bg = add_shape(slide, x, y, col_widths[j], row_height, bg_color, corner_radius=0.0)
            cell_bg.adjustments[0] = 0
            add_textbox(slide, x + Inches(0.15), y + Inches(0.07), col_widths[j] - Inches(0.2), row_height,
                        cell, font_size=11, color=DARK_TEXT if j == 0 else SUBTITLE_GRAY,
                        bold=(j == 0))
            x += col_widths[j]

    # Note
    add_textbox(slide, Inches(1.0), Inches(5.8), Inches(10), Inches(0.8),
                "All fees are transparently disclosed to users before transaction confirmation. "
                "Fee schedules will be filed with the Bank of Zambia as required under the "
                "National Payment Systems Act and related directives.",
                font_size=12, color=SUBTITLE_GRAY, line_spacing=1.4)

    # Revenue highlight
    card = add_shape(slide, Inches(1.0), Inches(6.5), Inches(10.9), Inches(0.7), RGBColor(0xEF, 0xF6, 0xFF), corner_radius=0.03)
    add_textbox(slide, Inches(1.3), Inches(6.6), Inches(10), Inches(0.5),
                "Monde collects all customer fees into a dedicated fee ledger account with full audit trail. "
                "Admin revenue withdrawal is a separate, controlled process.",
                font_size=12, color=MONDE_BLUE, bold=False)
    add_slide_number(slide, 9, TOTAL_SLIDES)

def slide_10_compliance_overview(prs):
    """Regulatory Compliance — NPSA 2007 Mapping"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "07  |  REGULATORY COMPLIANCE", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "NPSA 2007 — Compliance Assessment", font_size=32, color=DARK_TEXT, bold=True)

    add_textbox(slide, Inches(1.2), Inches(1.7), Inches(10), Inches(0.4),
                "Monde's compliance status against key sections of the National Payment Systems Act No. 1 of 2007",
                font_size=13, color=SUBTITLE_GRAY)

    # Compliance items
    items = [
        ("Section 3 — Application", "compliant",
         "Monde operates as a payment system business; Act applies"),
        ("Section 11 — BoZ Regulation", "partial",
         "Seeking designation; prepared for BoZ oversight"),
        ("Section 12 — Designation", "partial",
         "Application in preparation; this presentation is part of the process"),
        ("Section 13 — Restriction", "partial",
         "Not yet designated; operating in pre-launch/sandbox mode"),
        ("Section 27 — Returns", "partial",
         "System capable of generating reports; formal reporting structure to be agreed with BoZ"),
        ("Section 28 — Record Retention", "compliant",
         "All records stored electronically; 6+ year retention configured"),
        ("Section 29 — Access to Information", "compliant",
         "Full database access; can provide information within 14-day requirement"),
        ("Section 30 — Documents", "partial",
         "Governance structure being formalized; CEO/CFO signing authority to be established"),
        ("Section 33 — Dishonoured Cheques", "compliant",
         "Not applicable — Monde is a digital-only platform, no cheque processing"),
        ("Section 43 — Rules & Guidelines", "partial",
         "Ready to comply with any BoZ-prescribed rules, guidelines, or directives"),
    ]

    y = Inches(2.2)
    for label, status, detail in items:
        add_status_row(slide, Inches(1.0), y, Inches(11.0), label, status, detail)
        y += Inches(0.5)

    # Legend
    for i, (status, label, color) in enumerate([
        ("compliant", "Compliant / Ready", MONDE_GREEN),
        ("partial", "In Progress / Action Required", MONDE_GOLD),
        ("gap", "Gap / Not Yet Addressed", MONDE_RED)
    ]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0) + i * Inches(3.5),
                                      Inches(7.1), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide, Inches(1.2) + i * Inches(3.5), Inches(7.04), Inches(2.5), Inches(0.25),
                    label, font_size=10, color=SUBTITLE_GRAY)
    add_slide_number(slide, 10, TOTAL_SLIDES)

def slide_11_compliance_detail(prs):
    """Compliance Detail — What Monde Already Has"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MONDE_DARK)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "07B  |  COMPLIANCE DETAIL", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "What Monde Already Delivers", font_size=32, color=WHITE, bold=True)

    left_title = "Technical Compliance"
    left_items = [
        "Full transaction audit trail with timestamps and references",
        "Real-time balance tracking with database-level integrity",
        "Encrypted data transmission (HTTPS/TLS on all endpoints)",
        "Row-Level Security on all database tables",
        "JWT-based authentication with session management",
        "Input validation and sanitization on all user inputs",
        "Electronic record storage with 6+ year retention",
        "Separate fee ledger for revenue tracking and transparency",
    ]
    right_title = "Operational Compliance"
    right_items = [
        "KYC-ready architecture (phone-based identity verification)",
        "Transaction limits: min K1, max K50,000 per transaction",
        "Anti-fraud controls on agent operations",
        "Daily transaction caps and velocity checks",
        "Admin dashboard for oversight and monitoring",
        "Agent designation controlled by admin only",
        "Real-time notifications for all transactions",
        "Cross-provider fee transparency at point of transaction",
    ]

    # Left card
    card_l = add_shape(slide, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.0),
                       RGBColor(0x1A, 0x20, 0x35), corner_radius=0.03)
    card_l.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
    card_l.line.width = Pt(1)
    add_textbox(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.4),
                left_title, font_size=18, color=MONDE_GREEN, bold=True)
    for i, item in enumerate(left_items):
        add_textbox(slide, Inches(1.3), Inches(2.6) + i * Inches(0.35), Inches(5.0), Inches(0.3),
                    f"\u2713  {item}", font_size=11, color=MEDIUM_GRAY)

    # Right card
    card_r = add_shape(slide, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.0),
                       RGBColor(0x1A, 0x20, 0x35), corner_radius=0.03)
    card_r.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
    card_r.line.width = Pt(1)
    add_textbox(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.4),
                right_title, font_size=18, color=MONDE_BLUE, bold=True)
    for i, item in enumerate(right_items):
        add_textbox(slide, Inches(7.3), Inches(2.6) + i * Inches(0.35), Inches(5.0), Inches(0.3),
                    f"\u2713  {item}", font_size=11, color=MEDIUM_GRAY)
    add_slide_number(slide, 11, TOTAL_SLIDES)

def slide_12_compliance_roadmap(prs):
    """Compliance Roadmap — Steps Toward Full Designation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "08  |  COMPLIANCE ROADMAP", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Path to Full Designation Readiness", font_size=32, color=DARK_TEXT, bold=True)

    phases = [
        ("Phase 1: Immediate", "0–3 months", MONDE_BLUE, [
            "Submit formal designation application to BoZ (Section 12)",
            "Prepare and file all required documentation per BoZ prescribed forms",
            "Complete Directors' Questionnaire and Vital Statistics forms",
            "Engage legal counsel for compliance review",
            "Register company with PACRA (if not already done)",
        ]),
        ("Phase 2: Short-Term", "3–6 months", MONDE_GREEN, [
            "Implement formal KYC/CDD procedures aligned with FIC Act",
            "Establish AML/CFT compliance program and appoint MLRO",
            "File fee schedule with BoZ per Section 43 directives",
            "Set up formal returns/reporting structure per Section 27",
            "Establish governance structure with CEO/CFO signing authority",
        ]),
        ("Phase 3: Medium-Term", "6–12 months", MONDE_GOLD, [
            "Apply for BoZ Regulatory Sandbox (if applicable for testing)",
            "Implement customer complaint/grievance mechanism per Section 43",
            "Develop business continuity and disaster recovery plan",
            "Establish formal data protection policies per DPA 2021",
            "Build quarterly reporting framework for BoZ submissions",
        ]),
        ("Phase 4: Ongoing", "Continuous", RGBColor(0x8B, 0x5C, 0xF6), [
            "Regular compliance audits and self-assessments",
            "Respond to BoZ information requests within 14 days (Section 29)",
            "Maintain 6-year record retention (Section 28)",
            "Adapt to new BoZ directives and guidelines as issued",
            "Annual review of fee schedules and operational procedures",
        ]),
    ]

    card_w = Inches(5.5)
    card_h = Inches(2.35)
    for i, (title, timeline, color, items) in enumerate(phases):
        cx = Inches(0.8) + (i % 2) * (card_w + Inches(0.3))
        cy = Inches(1.9) + (i // 2) * (card_h + Inches(0.2))
        card = add_shape(slide, cx, cy, card_w, card_h, LIGHT_GRAY, corner_radius=0.03)
        # Colored top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(3.5), Inches(0.3),
                    title, font_size=14, color=DARK_TEXT, bold=True)
        add_textbox(slide, cx + card_w - Inches(1.5), cy + Inches(0.15), Inches(1.3), Inches(0.3),
                    timeline, font_size=11, color=color, bold=True, alignment=PP_ALIGN.RIGHT)
        for j, item in enumerate(items):
            add_textbox(slide, cx + Inches(0.3), cy + Inches(0.55) + j * Inches(0.32),
                        card_w - Inches(0.5), Inches(0.3),
                        f"\u2022  {item}", font_size=10, color=SUBTITLE_GRAY)
    add_slide_number(slide, 12, TOTAL_SLIDES)

def slide_13_risk_management(prs):
    """Risk Management & AML/CFT"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MONDE_DARK)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "09  |  RISK MANAGEMENT", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Risk Management & AML/CFT Framework", font_size=32, color=WHITE, bold=True)

    categories = [
        ("Operational Risk", MONDE_BLUE, [
            "Transaction amount limits (K1–K50,000) to contain exposure",
            "Real-time balance validation prevents overdrafts",
            "Atomic database transactions ensure consistency",
            "Automated error handling with graceful degradation",
            "Offline detection with user notification banner",
        ]),
        ("Fraud Prevention", MONDE_RED, [
            "Circular transaction detection (agent cash-in/cash-out within 24h)",
            "Daily velocity limits per customer and per agent",
            "Agent-only operations segregated from customer operations",
            "Self-service blocks prevent agents transacting with themselves",
            "Login rate limiting prevents brute-force attacks",
        ]),
        ("AML/CFT Readiness", MONDE_GREEN, [
            "Phone-based identity linked to MNO KYC records",
            "Full transaction audit trail for suspicious activity monitoring",
            "Planned: Formal MLRO appointment and STR filing procedures",
            "Planned: Integration with FIC reporting requirements",
            "Planned: Enhanced due diligence for high-value transactions",
        ]),
        ("Systemic Risk Mitigation", MONDE_GOLD, [
            "Funds held in regulated mobile money accounts (not pooled)",
            "Fee ledger separation ensures operational funds are distinct",
            "No credit extension — pure payment facilitation model",
            "Settlement via existing regulated payment infrastructure (Lipila/MNOs)",
            "Business continuity plan in development",
        ]),
    ]

    card_w = Inches(5.5)
    card_h = Inches(2.2)
    for i, (title, color, items) in enumerate(categories):
        cx = Inches(0.8) + (i % 2) * (card_w + Inches(0.3))
        cy = Inches(2.0) + (i // 2) * (card_h + Inches(0.2))
        card = add_shape(slide, cx, cy, card_w, card_h,
                        RGBColor(0x1A, 0x20, 0x35), corner_radius=0.03)
        card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
        card.line.width = Pt(1)
        # Left accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.06), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, cx + Inches(0.25), cy + Inches(0.12), card_w - Inches(0.4), Inches(0.3),
                    title, font_size=15, color=color, bold=True)
        for j, item in enumerate(items):
            add_textbox(slide, cx + Inches(0.35), cy + Inches(0.5) + j * Inches(0.3),
                        card_w - Inches(0.5), Inches(0.28),
                        f"\u2022  {item}", font_size=10, color=MEDIUM_GRAY)
    add_slide_number(slide, 13, TOTAL_SLIDES)

def slide_14_market(prs):
    """Market Opportunity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "10  |  MARKET OPPORTUNITY", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Zambia's Digital Payments Landscape", font_size=32, color=DARK_TEXT, bold=True)

    # Key stats
    stats = [
        ("20M+", "Population", "Growing youth demographic\ndriving mobile adoption", MONDE_BLUE),
        ("~60%", "Mobile Money\nPenetration", "Active mobile money accounts\ncontinue to grow rapidly", MONDE_GREEN),
        ("K200B+", "Annual Mobile\nMoney Value", "Transaction values growing\nyear-over-year", MONDE_GOLD),
        ("40%+", "Financially\nExcluded", "Significant population still\nlacking financial access", MONDE_RED),
    ]
    card_w = Inches(2.6)
    card_h = Inches(2.0)
    for i, (num, label, desc, color) in enumerate(stats):
        cx = Inches(0.9) + i * (card_w + Inches(0.2))
        cy = Inches(2.0)
        card = add_shape(slide, cx, cy, card_w, card_h, LIGHT_GRAY, corner_radius=0.04)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.2), card_w - Inches(0.3), Inches(0.5),
                    num, font_size=32, color=color, bold=True)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.8), card_w - Inches(0.3), Inches(0.4),
                    label, font_size=13, color=DARK_TEXT, bold=True)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(1.25), card_w - Inches(0.3), Inches(0.6),
                    desc, font_size=10, color=SUBTITLE_GRAY, line_spacing=1.3)

    # Why Monde
    add_textbox(slide, Inches(1.2), Inches(4.4), Inches(10), Inches(0.4),
                "Why Monde Matters", font_size=20, color=DARK_TEXT, bold=True)

    points = [
        ("Interoperability First",
         "Unlike single-provider wallets, Monde bridges ALL Zambian payment providers in one app, "
         "reducing friction and lowering costs for cross-provider transfers."),
        ("Agent-Driven Financial Inclusion",
         "Our agent cash-in/cash-out model provides physical touchpoints for the unbanked, "
         "enabling cash-to-digital conversion at the community level."),
        ("Regulatory Alignment",
         "Built from the ground up with BoZ regulatory requirements in mind. Seeking designation "
         "demonstrates our commitment to operating within the legal framework."),
    ]
    y = Inches(4.9)
    for title, desc in points:
        add_textbox(slide, Inches(1.4), y, Inches(2.5), Inches(0.3),
                    f"\u25B8  {title}", font_size=13, color=MONDE_BLUE, bold=True)
        add_textbox(slide, Inches(4.0), y, Inches(8), Inches(0.5),
                    desc, font_size=11, color=SUBTITLE_GRAY, line_spacing=1.3)
        y += Inches(0.65)
    add_slide_number(slide, 14, TOTAL_SLIDES)

def slide_15_timeline(prs):
    """Implementation Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "11  |  IMPLEMENTATION TIMELINE", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Phased Rollout Plan", font_size=32, color=DARK_TEXT, bold=True)

    phases = [
        ("Q1–Q2 2026", "Pre-Launch", MONDE_BLUE, [
            "Complete BoZ designation application",
            "Finalize legal and compliance framework",
            "Internal beta testing with controlled user group",
            "Establish agent onboarding process",
        ]),
        ("Q3 2026", "Sandbox / Pilot", MONDE_GREEN, [
            "Enter BoZ Regulatory Sandbox (if applicable)",
            "Pilot with limited agent network (50–100 agents)",
            "Gather user feedback and iterate",
            "File required reports and demonstrate compliance",
        ]),
        ("Q4 2026", "Controlled Launch", MONDE_GOLD, [
            "Obtain full BoZ designation",
            "Expand agent network nationwide",
            "Public launch with marketing campaign",
            "Establish BoZ reporting cadence",
        ]),
        ("2027+", "Scale & Expand", RGBColor(0x8B, 0x5C, 0xF6), [
            "Add merchant payment capabilities",
            "Expand to bill payments and utilities",
            "Explore cross-border payment corridors",
            "Continuous compliance and feature enhancement",
        ]),
    ]

    # Timeline bar
    bar_y = Inches(2.2)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.5), bar_y, Inches(10.3), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    bar.line.fill.background()

    # Phase nodes and cards
    card_w = Inches(2.5)
    card_h = Inches(3.5)
    for i, (time, name, color, items) in enumerate(phases):
        cx = Inches(1.2) + i * (card_w + Inches(0.35))
        # Node on timeline
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       cx + card_w/2 - Inches(0.12), bar_y - Inches(0.08),
                                       Inches(0.24), Inches(0.24))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        # Time label
        add_textbox(slide, cx, bar_y + Inches(0.3), card_w, Inches(0.3),
                    time, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # Phase name
        add_textbox(slide, cx, bar_y + Inches(0.6), card_w, Inches(0.3),
                    name, font_size=16, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
        # Items card
        card = add_shape(slide, cx, bar_y + Inches(1.0), card_w, Inches(2.2), LIGHT_GRAY, corner_radius=0.03)
        # Top bar
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, bar_y + Inches(1.0), card_w, Inches(0.05))
        top.fill.solid()
        top.fill.fore_color.rgb = color
        top.line.fill.background()
        for j, item in enumerate(items):
            add_textbox(slide, cx + Inches(0.15), bar_y + Inches(1.2) + j * Inches(0.42),
                        card_w - Inches(0.3), Inches(0.4),
                        f"\u2022  {item}", font_size=10, color=SUBTITLE_GRAY, line_spacing=1.2)
    add_slide_number(slide, 15, TOTAL_SLIDES)

def slide_16_team(prs):
    """Team & Governance (placeholder)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header_accent(slide)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "TEAM & GOVERNANCE", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "Leadership & Corporate Structure", font_size=32, color=DARK_TEXT, bold=True)

    # Placeholder cards for team members
    roles = [
        ("Founder / CEO", "Overall strategy, vision, and regulatory engagement with BoZ.\nResponsible for designation application and compliance oversight."),
        ("Chief Technology Officer", "Technology architecture, security, and platform development.\nEnsures system integrity, data protection, and technical compliance."),
        ("Chief Financial Officer", "Financial management, fee structures, and BoZ financial reporting.\nCo-signatory on official documents per Section 30."),
        ("Compliance Officer / MLRO", "AML/CFT compliance, KYC procedures, and regulatory reporting.\nLiaison with Financial Intelligence Centre (FIC)."),
    ]
    card_w = Inches(5.3)
    card_h = Inches(1.3)
    for i, (role, desc) in enumerate(roles):
        cx = Inches(1.0) + (i % 2) * (card_w + Inches(0.3))
        cy = Inches(2.0) + (i // 2) * (card_h + Inches(0.3))
        card = add_shape(slide, cx, cy, card_w, card_h, LIGHT_GRAY, corner_radius=0.03)
        # Avatar placeholder
        avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.2), cy + Inches(0.2),
                                         Inches(0.7), Inches(0.7))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = MONDE_BLUE
        avatar.line.fill.background()
        tf = avatar.text_frame
        p = tf.paragraphs[0]
        p.text = role[0]
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        add_textbox(slide, cx + Inches(1.1), cy + Inches(0.15), card_w - Inches(1.3), Inches(0.3),
                    role, font_size=14, color=DARK_TEXT, bold=True)
        add_textbox(slide, cx + Inches(1.1), cy + Inches(0.5), card_w - Inches(1.3), Inches(0.7),
                    desc, font_size=10, color=SUBTITLE_GRAY, line_spacing=1.3)

    # Governance note
    add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(1.5),
                "Corporate Governance Commitments:\n"
                "\u2022  Board of Directors with independent members to be established\n"
                "\u2022  Annual external audits by a BoZ-approved auditor\n"
                "\u2022  Internal compliance reviews on a quarterly basis\n"
                "\u2022  Segregation of duties between operations, finance, and compliance\n"
                "\u2022  Formal whistleblower and grievance procedures per Section 43(3)(g)\n"
                "\u2022  Company registration with PACRA and tax compliance with ZRA",
                font_size=12, color=SUBTITLE_GRAY, line_spacing=1.5)

    add_textbox(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(0.4),
                "Note: Specific names and CVs will be included in the formal application package alongside Directors' Questionnaires.",
                font_size=11, color=MEDIUM_GRAY)
    add_slide_number(slide, 16, TOTAL_SLIDES)

def slide_17_boz_sandbox(prs):
    """BoZ Regulatory Sandbox"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MONDE_DARK)
    add_textbox(slide, Inches(1.2), Inches(0.6), Inches(8), Inches(0.4),
                "REGULATORY SANDBOX", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.7),
                "BoZ Innovation Sandbox Readiness", font_size=32, color=WHITE, bold=True)

    add_textbox(slide, Inches(1.2), Inches(1.9), Inches(10), Inches(0.7),
                "Monde is prepared to participate in the Bank of Zambia's Regulatory Sandbox for innovative "
                "financial products and services, should BoZ recommend this pathway before full designation.",
                font_size=14, color=MEDIUM_GRAY, line_spacing=1.4)

    sandbox_items = [
        ("Innovative Product",
         "Monde's cross-provider interoperability and QR/NFC payment model represents a novel "
         "approach in the Zambian market, combining mobile money, banking, and agent networks "
         "in a single unified platform."),
        ("Consumer Benefit",
         "Lower transaction costs, faster transfers, and financial inclusion for unbanked "
         "populations through agent cash-in/cash-out services. Every transaction in 3 taps or fewer."),
        ("Testing Parameters",
         "We propose a controlled pilot with: limited geography (Lusaka initially), "
         "capped transaction volumes, restricted user base (500–1,000 users), "
         "and 50–100 agents for a 6-month testing period."),
        ("Monitoring & Reporting",
         "Weekly transaction reports, monthly compliance summaries, incident reporting "
         "within 24 hours, and full BoZ access to platform analytics and transaction data."),
        ("Exit Strategy",
         "Clear criteria for graduation to full designation or orderly wind-down. "
         "All customer funds protected through regulated MNO float accounts. "
         "No customer funds held directly by Monde."),
    ]

    card_w = Inches(10.8)
    y = Inches(2.8)
    for title, desc in sandbox_items:
        card = add_shape(slide, Inches(1.0), y, card_w, Inches(0.8),
                        RGBColor(0x1A, 0x20, 0x35), corner_radius=0.03)
        card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
        card.line.width = Pt(1)
        add_textbox(slide, Inches(1.3), y + Inches(0.08), Inches(2.5), Inches(0.25),
                    title, font_size=13, color=MONDE_GOLD, bold=True)
        add_textbox(slide, Inches(1.3), y + Inches(0.35), card_w - Inches(0.5), Inches(0.4),
                    desc, font_size=11, color=MEDIUM_GRAY, line_spacing=1.3)
        y += Inches(0.9)
    add_slide_number(slide, 17, TOTAL_SLIDES)

def slide_18_closing(prs):
    """Closing / Contact / Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MONDE_DARK)
    # Gold accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.0), Inches(2.0), Inches(0.08), Inches(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = MONDE_GOLD
    line.line.fill.background()

    add_textbox(slide, Inches(1.4), Inches(1.5), Inches(10), Inches(0.5),
                "NEXT STEPS", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.4), Inches(2.0), Inches(10), Inches(0.8),
                "Thank You", font_size=44, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.4), Inches(2.8), Inches(10), Inches(0.5),
                "Bank of Zambia — National Payment Systems Department",
                font_size=18, color=MEDIUM_GRAY)

    # Next steps
    steps = [
        "1.  Submit formal designation application with all required documents",
        "2.  Engage with BoZ on sandbox participation criteria (if recommended)",
        "3.  Complete Directors' Questionnaires and Vital Statistics forms",
        "4.  File proposed fee schedule for BoZ review and approval",
        "5.  Establish ongoing dialogue with NPS Department for guidance",
    ]
    y = Inches(3.7)
    for step in steps:
        add_textbox(slide, Inches(1.6), y, Inches(8), Inches(0.35),
                    step, font_size=14, color=MEDIUM_GRAY, line_spacing=1.3)
        y += Inches(0.4)

    # Contact info placeholder
    contact_card = add_shape(slide, Inches(1.0), Inches(5.8), Inches(5.5), Inches(1.2),
                             RGBColor(0x1A, 0x20, 0x35), corner_radius=0.04)
    contact_card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x55)
    contact_card.line.width = Pt(1)
    add_textbox(slide, Inches(1.3), Inches(5.9), Inches(5), Inches(0.3),
                "Contact Information", font_size=14, color=MONDE_GOLD, bold=True)
    add_textbox(slide, Inches(1.3), Inches(6.25), Inches(5), Inches(0.6),
                "Monde Payment Systems\n[Your Name, Title]\n[Email]  |  [Phone]",
                font_size=12, color=MEDIUM_GRAY, line_spacing=1.4)

    # Legal reference
    add_textbox(slide, Inches(7.0), Inches(5.8), Inches(5.5), Inches(1.2),
                "This application is made pursuant to:\n\n"
                "\u2022  National Payment Systems Act No. 1 of 2007, Section 12\n"
                "\u2022  BoZ Requirements for Setting Up a Payment System Business\n"
                "\u2022  BoZ Requirements for Designation of a Payment System\n"
                "\u2022  BoZ Payment Service Provider Fees and Charges Guidelines",
                font_size=11, color=MEDIUM_GRAY, line_spacing=1.4)

    # Bottom bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MONDE_BLUE
    bar.line.fill.background()
    add_slide_number(slide, 18, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_exec_summary(prs)
    slide_04_problem(prs)
    slide_05_solution(prs)
    slide_06_tech_security(prs)
    slide_07_data_protection(prs)
    slide_08_agent_network(prs)
    slide_09_business_model(prs)
    slide_10_compliance_overview(prs)
    slide_11_compliance_detail(prs)
    slide_12_compliance_roadmap(prs)
    slide_13_risk_management(prs)
    slide_14_market(prs)
    slide_15_timeline(prs)
    slide_16_team(prs)
    slide_17_boz_sandbox(prs)
    slide_18_closing(prs)

    # Save
    out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    out_path = os.path.join(out_dir, "Monde_BoZ_Presentation.pptx")
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {TOTAL_SLIDES}")

if __name__ == "__main__":
    main()
